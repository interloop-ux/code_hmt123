from pptx import Presentation
from pptx.util import Inches, Pt

def create_project_ppt():
    prs = Presentation()

    # --- SLIDE 1: TITLE SLIDE ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "iOS Health Tracker Project"
    slide.placeholders[1].text = "Developed in Python & VS Code\nStep Tracking & HealthKit Integration"

    # --- SLIDE 2: PROJECT GOAL ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Project Overview"
    tf = slide.placeholders[1].text_frame
    tf.text = "Goal: Create a native iOS application to track daily steps."
    p = tf.add_paragraph()
    p.text = "Key Innovation: Using Python as the primary development language instead of Swift."

    # --- SLIDE 3: TECH STACK ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "The Tech Stack"
    tf = slide.placeholders[1].text_frame
    tf.text = "Editor: VS Code"
    tf.add_paragraph().text = "Language: Python 3.x"
    tf.add_paragraph().text = "Framework: Kivy (for the iOS User Interface)"
    tf.add_paragraph().text = "Bridge: Pyobjus (to talk to Apple's Hardware)"

    # --- SLIDE 4: THE STEP COUNTER LOGIC ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Core Functionality: Step Tracking"
    tf = slide.placeholders[1].text_frame
    tf.text = "Connects to Apple HealthKit API."
    tf.add_paragraph().text = "Requests user permission for 'HKQuantityTypeIdentifierStepCount'."
    tf.add_paragraph().text = "Calculates total steps from the start of the current day."

    # --- SLIDE 5: DEVELOPMENT WORKFLOW ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Development Workflow"
    tf = slide.placeholders[1].text_frame
    tf.text = "1. Code the UI & Logic in Python (VS Code)."
    tf.add_paragraph().text = "2. Use Kivy-ios toolchain to package for iOS."
    tf.add_paragraph().text = "3. Add Privacy Descriptions to Info.plist."
    tf.add_paragraph().text = "4. Deploy to iPhone via Xcode."

    # --- SLIDE 6: CHALLENGES & SOLUTIONS ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Challenges & Solutions"
    tf = slide.placeholders[1].text_frame
    tf.text = "Challenge: Python isn't native to iOS."
    tf.add_paragraph().text = "Solution: Embedding a Python interpreter within the app bundle."
    tf.add_paragraph().text = "Challenge: Hardware access."
    tf.add_paragraph().text = "Solution: Using Objective-C bridges for HealthKit."

    # --- SLIDE 7: CONCLUSION ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Conclusion"
    tf = slide.placeholders[1].text_frame
    tf.text = "This project demonstrates that Python is a viable tool for health-tech mobile development, allowing for rapid prototyping and cross-platform potential."

    # Save the file
    prs.save('My_Health_App_Project.pptx')
    print("Presentation created successfully as 'My_Health_App_Project.pptx'")

if __name__ == "__main__":
    create_project_ppt()